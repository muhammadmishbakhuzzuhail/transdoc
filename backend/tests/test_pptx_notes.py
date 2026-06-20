# SPDX-License-Identifier: AGPL-3.0-only
# Copyright (C) 2026 Muhammad Mishbakhuz Zuhail
"""PPTX speaker notes are extracted (were silently dropped) and round-trip through the renderer."""

from __future__ import annotations

import pytest

from transdoc.config import Config

pytest.importorskip("pptx")


def _deck(path):
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Slide Title"
    slide.notes_slide.notes_text_frame.text = "Hidden speaker note"
    prs.save(str(path))


def test_speaker_notes_extracted_and_rendered(tmp_path):
    from transdoc.extract.pptx import extract
    from transdoc.regenerate.pptx_out import render
    src = tmp_path / "d.pptx"
    _deck(src)
    doc = extract(str(src), Config(target_lang="id"))
    note = next((b for b in doc.blocks if "speaker note" in (b.text or "").lower()), None)
    assert note is not None and "notes" in note.id            # notes captured with a notes id
    note.translated = "Catatan pembicara"
    out = tmp_path / "out.pptx"
    render(doc, Config(target_lang="id"), str(out))
    from pptx import Presentation
    prs = Presentation(str(out))
    notes_text = prs.slides[0].notes_slide.notes_text_frame.text
    assert "Catatan pembicara" in notes_text                  # translation written back to notes
