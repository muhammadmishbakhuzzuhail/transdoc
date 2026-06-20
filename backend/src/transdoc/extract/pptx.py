# SPDX-License-Identifier: AGPL-3.0-only
# Copyright (C) 2026 Muhammad Mishbakhuz Zuhail
"""PowerPoint (.pptx) extraction via python-pptx.

One IR block per text-bearing paragraph (in shapes and table cells), with a stable id that
encodes its location (slide/shape/paragraph). The renderer reopens the source deck and swaps
text by that id, so every visual property — theme, position, fonts, animations — is kept.
"""

from __future__ import annotations

from ..config import Config
from ..ir import Block, BlockType, Confidence, Document
from .base import reflow_order


def _para_text(para) -> str:
    return "".join(run.text for run in para.runs)


def _iter_shape_paras(shape, prefix):
    """Yield (id, paragraph) for a shape, recursing into GROUP shapes (whose text was previously
    dropped). Stable ids encode the shape path so extract + render walk identically."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for gi, sub in enumerate(shape.shapes):
            yield from _iter_shape_paras(sub, f"{prefix}_g{gi}")
        return
    if shape.has_text_frame:
        for pi, para in enumerate(shape.text_frame.paragraphs):
            yield f"{prefix}_p{pi}", para
    elif shape.has_table:
        for ri, row in enumerate(shape.table.rows):
            for ci, cell in enumerate(row.cells):
                for pi, para in enumerate(cell.text_frame.paragraphs):
                    yield f"{prefix}_t{ri}_{ci}_p{pi}", para


def iter_text_paras(prs):
    """The single walk over a deck's translatable paragraphs — shapes (incl. grouped), table
    cells, AND speaker notes — yielding (id, paragraph, slide_index). Used by BOTH the extractor
    and the renderer so they can never drift out of sync."""
    for si, slide in enumerate(prs.slides):
        for shi, shape in enumerate(slide.shapes):
            yield from ((bid, para, si) for bid, para in _iter_shape_paras(shape, f"s{si}_sh{shi}"))
        if slide.has_notes_slide:                       # speaker notes were silently dropped
            nf = slide.notes_slide.notes_text_frame
            if nf is not None:
                for pi, para in enumerate(nf.paragraphs):
                    yield f"s{si}_notes_p{pi}", para, si


def extract(path: str, cfg: Config) -> Document:
    from pptx import Presentation

    try:
        prs = Presentation(path)
    except Exception as e:
        raise ValueError(f"unreadable or corrupt PPTX: {e}") from e
    out = Document(source_path=path,
                   mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    for bid, para, si in iter_text_paras(prs):
        t = _para_text(para)
        if t.strip():
            out.blocks.append(Block(id=bid, type=BlockType.PARAGRAPH, page=si, text=t,
                                    confidence=Confidence(source="digital")))
    reflow_order(out)
    return out
